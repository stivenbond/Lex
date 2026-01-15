import docx
import pptx
import os

class DocumentParser:
    @staticmethod
    def parse_file(file_path):
        """
        Parses a file (.docx or .pptx) and returns a list of text blocks.
        Result format: [{"type": "paragraph", "content": "text..."}, ...]
        """
        ext = os.path.splitext(file_path)[1].lower()
        
        if ext == ".docx":
            return DocumentParser._parse_docx(file_path)
        elif ext == ".pptx":
            return DocumentParser._parse_pptx(file_path)
        else:
            raise ValueError(f"Unsupported file format: {ext}")

    @staticmethod
    def _parse_docx(file_path):
        blocks = []
        doc = docx.Document(file_path)
        
        for para in doc.paragraphs:
            if para.text.strip():
                blocks.append({
                    "type": "paragraph",
                    "content": para.text.strip()
                })
        return blocks

    @staticmethod
    def _parse_pptx(file_path):
        blocks = []
        prs = pptx.Presentation(file_path)
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    if shape.text_frame.text.strip():
                        blocks.append({
                            "type": "paragraph",
                            "content": shape.text_frame.text.strip()
                        })
        return blocks
