import docx
from pptx import Presentation
import sqlite3

def parse_docx_lessons(file_path):
    doc = docx.Document(file_path)
    blocks = []
    for para in doc.paragraphs:
        if para.text.strip():
            blocks.append(para.text.strip())
    return blocks

def parse_pptx_lessons(file_path):
    prs = Presentation(file_path)
    blocks = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                blocks.append(shape.text.strip())
    return blocks

def save_lesson_content(topic_id, blocks):
    conn = sqlite3.connect("lex_database.db")
    cursor = conn.cursor()
    
    # Clear existing blocks for this topic if any
    cursor.execute("DELETE FROM text_blocks WHERE topic_id = ?", (topic_id,))
    
    for i, content in enumerate(blocks):
        cursor.execute("INSERT INTO text_blocks (topic_id, content, order_index) VALUES (?, ?, ?)",
                       (topic_id, content, i))
    
    conn.commit()
    conn.close()
